import asyncio

from pptx import Presentation  # type: ignore

from content_core.logging import logger


async def extract_pptx_content(file_path):
    """Extract content from PPTX file"""

    def _extract():
        try:
            prs = Presentation(file_path)
            content = []

            for slide_number, slide in enumerate(prs.slides, 1):
                content.append(f"\n# Slide {slide_number}\n")

                # Extract title
                if slide.shapes.title:
                    content.append(f"## {slide.shapes.title.text}\n")

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if (
                            shape != slide.shapes.title
                        ):  # Skip title as it's already added
                            content.append(shape.text.strip())

            return "\n\n".join(content)

        except Exception as e:
            logger.error(f"Failed to extract PPTX content: {e}")
            return None

    return await asyncio.get_event_loop().run_in_executor(None, _extract)


async def get_pptx_info(file_path):
    """Get PPTX metadata and content"""

    def _get_pptx_metadata_sync(file_path):
        """Synchronous helper to extract metadata using python-pptx."""
        try:
            prs = Presentation(file_path)
            props = {
                "slide_count": len(prs.slides),
                "title": "",  # PowerPoint doesn't have built-in metadata like Word
            }
            stats = {
                "slide_count": len(prs.slides),
                "shape_count": sum(len(slide.shapes) for slide in prs.slides),
                "text_frame_count": sum(
                    sum(1 for shape in slide.shapes if hasattr(shape, "text"))
                    for slide in prs.slides
                ),
            }
            return {"metadata": props, "statistics": stats}
        except Exception as e:
            logger.error(f"Failed to get PPTX metadata: {e}")
            return None

    try:
        # Run blocking python-pptx operations in executor
        metadata_info = await asyncio.get_event_loop().run_in_executor(
            None, _get_pptx_metadata_sync, file_path
        )

        # Await the async content extraction directly
        content = await extract_pptx_content(file_path)

        if metadata_info:
            # Combine results
            return {**metadata_info, "content": content}
        else:
            # Fallback if metadata extraction failed
            return {"metadata": {}, "statistics": {}, "content": content}

    except Exception as e:
        logger.error(f"Failed to get PPTX info: {e}")
        return None
